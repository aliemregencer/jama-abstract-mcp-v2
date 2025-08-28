# app.py dosyasÄ±nÄ±n tamamÄ±nÄ± bu gÃ¼ncel versiyonla deÄŸiÅŸtir

import json
import re
import os
import base64
from pprint import pprint
from bs4 import BeautifulSoup
from selenium import webdriver
from selenium.webdriver.chrome.service import Service as ChromeService
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.chrome.options import Options
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import requests
import time
from datetime import datetime

# --- Ä°KON EÅLEÅTÄ°RME HARÄ°TASI (DeÄŸiÅŸiklik yok) ---
ICON_MAP = {
    "cardiology.png": ["heart", "cardiac", "cardiology", "myocardial", "arrhythmia", "heart failure"],
    "neurology.png": ["brain", "neuro", "neurology", "stroke", "alzheimer", "parkinson", "epilepsy"],
    "oncology.png": ["cancer", "oncology", "tumor", "chemotherapy", "carcinoma"],
    "public_health.png": ["population", "public health", "mortality", "epidemiology", "opioid", "addiction"],
    "genetics.png": ["gene", "genetic", "dna", "genome", "genomics"],
}

def select_thematic_icon(article_title, article_keywords):
    # Bu fonksiyonda deÄŸiÅŸiklik yok
    search_text = (article_title.lower() + ' ' + ' '.join(article_keywords)).lower()
    for icon_file, keywords in ICON_MAP.items():
        for keyword in keywords:
            if keyword in search_text:
                return f"icons/{icon_file}"
    return "icons/default.png"

def parse_jama_article(url):
    # GÃœNCELLEME: Container ortamÄ±nda Ã§alÄ±ÅŸacak ÅŸekilde Selenium konfigÃ¼rasyonu
    html_content = None
    
    # Ã–nce requests ile deneyelim (daha hÄ±zlÄ± ve gÃ¼venilir)
    print("ğŸ“¡ Requests ile sayfa yÃ¼kleniyor...")
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.5',
            'Accept-Encoding': 'gzip, deflate',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
        }
        
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        html_content = response.text
        
        # Sayfa yÃ¼klendi mi kontrol et
        if "jamanetwork" in html_content.lower() and len(html_content) > 1000:
            print("âœ… Requests ile sayfa baÅŸarÄ±yla yÃ¼klendi")
        else:
            print("âš ï¸ Requests ile sayfa yÃ¼klendi ama iÃ§erik eksik, Selenium deneniyor...")
            html_content = None
            
    except Exception as e:
        print(f"âš ï¸ Requests ile yÃ¼kleme baÅŸarÄ±sÄ±z: {e}")
        html_content = None
    
    # Requests baÅŸarÄ±sÄ±z olursa Selenium'u dene
    if not html_content:
        print("ğŸ”„ Selenium ile sayfa yÃ¼kleniyor...")
        try:
            # Container ortamÄ±nda Ã§alÄ±ÅŸacak Chrome options
            chrome_options = Options()
            chrome_options.add_argument("--headless")
            chrome_options.add_argument("--no-sandbox")
            chrome_options.add_argument("--disable-dev-shm-usage")
            chrome_options.add_argument("--disable-gpu")
            chrome_options.add_argument("--disable-extensions")
            chrome_options.add_argument("--disable-plugins")
            chrome_options.add_argument("--disable-images")
            chrome_options.add_argument("--disable-javascript")  # JavaScript'i devre dÄ±ÅŸÄ± bÄ±rak
            chrome_options.add_argument("--window-size=1920,1080")
            chrome_options.add_argument("--log-level=3")
            chrome_options.add_argument("--silent")
            chrome_options.add_argument("--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36")
            
            # Container ortamÄ±nda Chrome driver kurulumu
            try:
                # Ã–nce sistem Chrome'u kullanmayÄ± dene
                service = ChromeService()
                driver = webdriver.Chrome(service=service, options=chrome_options)
                print("âœ… Sistem Chrome driver kullanÄ±lÄ±yor")
            except:
                # Sistem Chrome yoksa ChromeDriverManager kullan
                try:
                    # Container ortamÄ±nda Chrome driver kurulumu
                    import os
                    chrome_bin = os.getenv('CHROME_BIN', '/usr/bin/google-chrome')
                    if os.path.exists(chrome_bin):
                        chrome_options.binary_location = chrome_bin
                        print(f"âœ… Chrome binary bulundu: {chrome_bin}")
                    
                    # ChromeDriverManager'Ä± container ortamÄ±nda Ã§alÄ±ÅŸacak ÅŸekilde ayarla
                    os.environ['WDM_LOG_LEVEL'] = '0'  # Log seviyesini dÃ¼ÅŸÃ¼r
                    os.environ['WDM_LOCAL'] = '1'      # Yerel cache kullan
                    
                    service = ChromeService(ChromeDriverManager().install())
                    driver = webdriver.Chrome(service=service, options=chrome_options)
                    print("âœ… ChromeDriverManager ile Chrome driver kuruldu")
                except Exception as driver_error:
                    print(f"âŒ Chrome driver kurulumu baÅŸarÄ±sÄ±z: {driver_error}")
                    # Son Ã§are: requests ile tekrar dene
                    print("ğŸ”„ Son Ã§are: Requests ile tekrar deneniyor...")
                    try:
                        response = requests.get(url, headers=headers, timeout=60)
                        response.raise_for_status()
                        html_content = response.text
                        if "jamanetwork" in html_content.lower() and len(html_content) > 1000:
                            print("âœ… Son Ã§are requests baÅŸarÄ±lÄ±")
                        else:
                            raise Exception("HTML iÃ§erik yetersiz")
                    except Exception as final_error:
                        error_message = f"TÃ¼m scraping yÃ¶ntemleri baÅŸarÄ±sÄ±z: {final_error}"
                        print(error_message)
                        return None, error_message
            
            # Selenium ile sayfa yÃ¼kle
            if 'driver' in locals():
                driver.set_page_load_timeout(60)
                driver.get(url)
                
                # Sayfa yÃ¼klenene kadar bekle
                time.sleep(5)
                
                # JavaScript'i etkinleÅŸtir ve tekrar yÃ¼kle
                if "jamanetwork" not in driver.page_source.lower():
                    print("ğŸ”„ JavaScript ile tekrar yÃ¼kleniyor...")
                    chrome_options.remove_argument("--disable-javascript")
                    driver.quit()
                    
                    service = ChromeService()
                    driver = webdriver.Chrome(service=service, options=chrome_options)
                    driver.set_page_load_timeout(60)
                    driver.get(url)
                    time.sleep(5)
                
                html_content = driver.page_source
                
        except Exception as e:
            error_message = f"Selenium ile sayfa yÃ¼klenirken bir hata oluÅŸtu: {str(e)}"
            print(error_message)
            
            # Son Ã§are: requests ile tekrar dene
            try:
                print("ğŸ”„ Son Ã§are: Requests ile tekrar deneniyor...")
                response = requests.get(url, headers=headers, timeout=60)
                response.raise_for_status()
                html_content = response.text
                if "jamanetwork" in html_content.lower() and len(html_content) > 1000:
                    print("âœ… Son Ã§are requests baÅŸarÄ±lÄ±")
                else:
                    return None, error_message
            except Exception as final_error:
                return None, f"TÃ¼m scraping yÃ¶ntemleri baÅŸarÄ±sÄ±z. Son hata: {final_error}"
        finally:
            if 'driver' in locals():
                driver.quit()

    if not html_content:
        error_message = "HTML iÃ§erik alÄ±namadÄ± (sayfa boÅŸ geldi)."
        print(error_message)
        return None, error_message

    print("âœ… HTML iÃ§erik baÅŸarÄ±yla alÄ±ndÄ±, parsing baÅŸlÄ±yor...")
    
    # ... (Geri kalan parsing kodu aynÄ±, sadece en sonda return deÄŸerini gÃ¼ncelleyeceÄŸiz)
    soup = BeautifulSoup(html_content, 'html.parser')
    article_data = {
        'url': url, 'title': None, 'authors': [], 'publication_date': None,
        'doi': None, 'key_points': {}, 'abstract': {}, 'keywords': [],
        'full_text': {}, 'references': []
    }
    
    # JSON-LD verilerini parse et
    try:
        for script_tag in soup.find_all("script", {"type": "application/ld+json"}):
            if script_tag.string:
                json_content = json.loads(script_tag.string)
                if json_content.get("@type") == "MedicalScholarlyArticle":
                    article_data["publication_date"] = json_content.get("datePublished")
                    article_data["keywords"] = json_content.get("keyWords", "").split(", ")
                    break
    except (json.JSONDecodeError, AttributeError, TypeError):
        print("JSON-LD verisi ayrÄ±ÅŸtÄ±rÄ±lÄ±rken bir sorun oluÅŸtu, devam ediliyor.")
    
    # BaÅŸlÄ±k
    title_tag = soup.find('h1', class_='meta-article-title')
    if title_tag:
        article_data['title'] = title_tag.get_text(strip=True)
    else:
        # GÃœNCELLEME: Parsing baÅŸarÄ±sÄ±z olursa da hata dÃ¶ndÃ¼r
        error_message = "Makale baÅŸlÄ±ÄŸÄ± (title) bulunamadÄ±. Sayfa yapÄ±sÄ± deÄŸiÅŸmiÅŸ olabilir."
        print(error_message)
        return None, error_message
    
    # Yazarlar
    authors_limited = soup.select(".meta-authors--limited .wi-fullname")
    authors_remaining = soup.select(".meta-authors--remaining .wi-fullname")
    all_authors = authors_limited + authors_remaining
    author_affiliations_map = {}
    affiliations_list = soup.select(".meta-author-affiliations li")
    for li in affiliations_list:
        sup_tag = li.find("sup")
        if sup_tag:
            sup_num = sup_tag.get_text(strip=True)
            affiliation_text_div = li.find("div", class_="meta-author-name")
            if affiliation_text_div:
                affiliation_text = (
                    affiliation_text_div.get_text(strip=True)
                    .replace(sup_num, "", 1)
                    .strip()
                )
                author_affiliations_map[sup_num] = affiliation_text
    for author in all_authors:
        name_tag = author.find("a") if author.find("a") else author
        sup_tags = name_tag.find_all("sup")
        aff_keys = [sup.get_text(strip=True) for sup in sup_tags]
        for sup in name_tag.find_all("sup"):
            sup.decompose()
        name = name_tag.get_text(strip=True).replace(",", "")
        author_info = {
            "name": name,
            "affiliations": [
                author_affiliations_map.get(key)
                for key in aff_keys
                if author_affiliations_map.get(key)
            ],
        }
        article_data["authors"].append(author_info)
    
    # DOI
    doi_tag = soup.find("span", class_="meta-citation-doi")
    if doi_tag:
        article_data["doi"] = doi_tag.get_text(strip=True).replace("doi:", "")
    
    # Key Points
    key_points_section = soup.find("span", class_="heading-text", string="Key Points")
    if key_points_section:
        current_element = key_points_section.find_next_sibling("p")
        while current_element and current_element.name == "p":
            strong_tag = current_element.find("strong")
            if strong_tag:
                key = strong_tag.get_text(strip=True).lower()
                value = (
                    current_element.find("span").get_text(strip=True)
                    if current_element.find("span")
                    else ""
                )
                article_data["key_points"][key] = value
            current_element = current_element.find_next_sibling()
    
    # Abstract
    abstract_section = soup.find("div", id="AbstractSection")
    if abstract_section:
        for p_tag in abstract_section.find_all("p"):
            strong_tag = p_tag.find("strong")
            if strong_tag:
                key = strong_tag.get_text(strip=True).lower().replace(":", "")
                span_tag = p_tag.find("span")
                if span_tag:
                    article_data["abstract"][key] = span_tag.get_text(strip=True)
    
    # Full Text
    full_text_div = soup.find("div", class_="article-full-text")
    if full_text_div:
        first_header = full_text_div.find("div", class_="section-type-section")
        if first_header:
            current_section_title = "introduction"
            article_data["full_text"][current_section_title] = ""
            for element in first_header.find_all_previous("p", class_="para"):
                article_data["full_text"][current_section_title] = (
                    element.get_text(" ", strip=True)
                    + "\n"
                    + article_data["full_text"][current_section_title]
                )

            for element in first_header.find_next_siblings():
                if element.name == "div" and "section-type-section" in element.get(
                    "class", []
                ):
                    current_section_title = element.get_text(strip=True).lower()
                    article_data["full_text"][current_section_title] = ""
                elif element.name == "p" and "para" in element.get("class", []):
                    if current_section_title:
                        article_data["full_text"][current_section_title] += (
                            element.get_text(" ", strip=True) + "\n"
                        )
    
    # References
    references_div = soup.find("div", class_="references")
    if references_div:
        for ref in references_div.find_all("div", class_="reference"):
            ref_content = ref.find("div", class_="reference-content")
            if ref_content:
                article_data["references"].append(ref_content.get_text(" ", strip=True))
    
    return article_data, None # BaÅŸarÄ±lÄ± durumda (data, None) dÃ¶ndÃ¼r

def create_presentation(data, icon_path):
    """
    Verilen verilerle jama_va.pptx template'ini kullanarak bir PowerPoint sunumu oluÅŸturur.
    Template'deki ÅŸekil isimlerine gÃ¶re iÃ§erik yerleÅŸtirir.
    """
    # Template dosyasÄ±nÄ± aÃ§
    template_path = "templates/jama_va.pptx"
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template dosyasÄ± bulunamadÄ±: {template_path}")
    
    prs = Presentation(template_path)
    
    # Template'deki ÅŸekilleri bul ve iÃ§erikle doldur
    for slide in prs.slides:
        for shape in slide.shapes:
            shape_name = shape.name if hasattr(shape, 'name') else ""
            
            # Åekil ismine gÃ¶re iÃ§erik yerleÅŸtir
            if shape_name == "title":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = data.get("title", "BaÅŸlÄ±k BulunamadÄ±")
                    p.font.bold = True
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0xED, 0x09, 0x73)
            
            elif shape_name == "population_subtitle":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = "POPULATION"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(0xED, 0x09, 0x73)
            
            elif shape_name == "population_description":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    # Makale iÃ§eriÄŸinden popÃ¼lasyon bilgisini Ã§Ä±kar
                    population_text = extract_population_info(data)
                    p.text = population_text
                    p.font.size = Pt(10)
                    # Metin kutusuna sÄ±ÄŸdÄ±r
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    shape.text_frame.word_wrap = True
            
            elif shape_name == "intervention_subtitle":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = "INTERVENTION"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(0xED, 0x09, 0x73)
            
            elif shape_name == "intervention_description":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    # Makale iÃ§eriÄŸinden mÃ¼dahale bilgisini Ã§Ä±kar
                    intervention_text = extract_intervention_info(data)
                    p.text = intervention_text
                    p.font.size = Pt(10)
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    shape.text_frame.word_wrap = True
            
            elif shape_name == "settings_locations_description":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    # Makale iÃ§eriÄŸinden ayarlar ve konum bilgisini Ã§Ä±kar
                    settings_text = extract_settings_info(data)
                    p.text = settings_text
                    p.font.size = Pt(10)
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    shape.text_frame.word_wrap = True
            
            elif shape_name == "primary_outcome_description":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    # Makale iÃ§eriÄŸinden birincil sonuÃ§ bilgisini Ã§Ä±kar
                    outcome_text = extract_primary_outcome_info(data)
                    p.text = outcome_text
                    p.font.size = Pt(10)
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    shape.text_frame.word_wrap = True
            
            elif shape_name == "findings_description_1":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    # Makale iÃ§eriÄŸinden bulgular bilgisini Ã§Ä±kar
                    findings_text = extract_findings_info(data)
                    # BulgularÄ± iki parÃ§aya bÃ¶l
                    if findings_text and len(findings_text) > 200:
                        words = findings_text.split()
                        mid_point = len(words) // 2
                        p.text = " ".join(words[:mid_point])
                    else:
                        p.text = findings_text
                    p.font.size = Pt(10)
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    shape.text_frame.word_wrap = True
            
            elif shape_name == "findings_description_2":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    # Makale iÃ§eriÄŸinden bulgular bilgisini Ã§Ä±kar
                    findings_text = extract_findings_info(data)
                    # BulgularÄ± iki parÃ§aya bÃ¶l
                    if findings_text and len(findings_text) > 200:
                        words = findings_text.split()
                        mid_point = len(words) // 2
                        p.text = " ".join(words[mid_point:])
                    else:
                        p.text = ""
                    p.font.size = Pt(10)
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    shape.text_frame.word_wrap = True
            
            elif shape_name == "footer_citation":
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    first_author = data["authors"][0]["name"] if data["authors"] else "Yazar Yok"
                    citation = f"{first_author.split(',')[0]} L, et al. {data.get('title', 'BaÅŸlÄ±k')}; a randomized clinical trial. JAMA Netw Open. {data.get('publication_date', 'Tarih Yok')} doi:{data.get('doi', 'DOI Yok')}"
                    p.text = citation
                    p.font.size = Pt(8)
                    p.font.color.rgb = RGBColor(128, 128, 128)
            
            # Metin kutularÄ±nÄ± da kontrol et - ana baÅŸlÄ±klarÄ± deÄŸiÅŸtirme
            elif "Metin kutusu" in shape_name:
                if hasattr(shape, 'text_frame'):
                    # Metin kutusunun konumuna gÃ¶re iÃ§erik yerleÅŸtir
                    left = shape.left
                    top = shape.top
                    
                    # Sol taraftaki metin kutularÄ± (popÃ¼lasyon alanÄ±)
                    if left < Inches(3):
                        if "16" in shape_name or "17" in shape_name:
                            # Bu metin kutularÄ± ana baÅŸlÄ±k, deÄŸiÅŸtirme
                            continue
                    
                    # SaÄŸ taraftaki metin kutularÄ± (bulgular alanÄ±)
                    elif left > Inches(3):
                        if "20" in shape_name or "21" in shape_name or "22" in shape_name:
                            # Bu metin kutularÄ± ana baÅŸlÄ±k, deÄŸiÅŸtirme
                            continue

    # DosyayÄ± kaydet
    filename = "JAMA_VA_Abstract.pptx"
    prs.save(filename)
    print(f"\nVA formatÄ±nda sunum baÅŸarÄ±yla oluÅŸturuldu ve kaydedildi: {filename}")
    return filename

def upload_to_github_release(
    filename: str,
    title: str,
    repo_full_name: str,
    github_token: str,
) -> tuple[str | None, str | None]:
    """
    Verilen `repo_full_name` (Ã¶rn. "kullaniciadi/repoadi") ve `github_token` ile GitHub'da
    `latest-abstract` etiketiyle bir Release oluÅŸturur, varsa eskisini siler ve `filename`
    dosyasÄ±nÄ± bu release'e asset olarak yÃ¼kler. BaÅŸarÄ±lÄ±ysa herkese aÃ§Ä±k indirme linkini dÃ¶ndÃ¼rÃ¼r.
    """
    try:
        if not repo_full_name or "/" not in repo_full_name:
            msg = "GeÃ§ersiz repo formatÄ±. 'kullaniciadi/repoadi' ÅŸeklinde olmalÄ±."
            print(msg)
            return None, msg
        if not github_token:
            msg = "GitHub token gerekli, iÅŸlem iptal edildi."
            print(msg)
            return None, msg

        repo_owner, repo_name = repo_full_name.split("/", 1)
        release_tag = "latest-abstract"
        safe_title = title or "JAMA Abstract"
        release_name = f"JAMA Abstract - {safe_title[:70]}"  # daha kÄ±sa

        # API headers
        headers_json = {
            "Authorization": f"Bearer {github_token}",
            "Accept": "application/vnd.github+json",
        }

        api_base = f"https://api.github.com/repos/{repo_owner}/{repo_name}"

        # Token ve repo eriÅŸimini doÄŸrula
        repo_check = requests.get(api_base, headers=headers_json)
        if repo_check.status_code != 200:
            msg = (
                f"Repo eriÅŸimi baÅŸarÄ±sÄ±z: {repo_check.status_code} {repo_check.text}. "
                "Repo adÄ±nÄ± ve token izinlerini kontrol edin."
            )
            print(msg)
            return None, msg

        # Mevcut release'i kontrol et ve varsa sil
        print("Mevcut release kontrol ediliyor...")
        response = requests.get(f"{api_base}/releases/tags/{release_tag}", headers=headers_json)
        if response.status_code == 200:
            release = response.json()
            release_id = release["id"]
            # Asset'leri sil
            assets_resp = requests.get(f"{api_base}/releases/{release_id}/assets", headers=headers_json)
            if assets_resp.status_code == 200:
                for asset in assets_resp.json():
                    requests.delete(f"{api_base}/releases/assets/{asset['id']}", headers=headers_json)
            # Release'i sil
            requests.delete(f"{api_base}/releases/{release_id}", headers=headers_json)
            # Tag'i de sil (aksi halde aynÄ± tag ile oluÅŸturma baÅŸarÄ±sÄ±z olabilir)
            requests.delete(f"{api_base}/git/refs/tags/{release_tag}", headers=headers_json)
            print("Eski release ve etiketi silindi.")
        elif response.status_code not in (200, 404):
            msg = f"Release kontrolÃ¼nde hata: {response.status_code} {response.text}"
            print(msg)
            return None, msg

        # Yeni release oluÅŸtur
        print("Yeni release oluÅŸturuluyor...")
        release_data = {
            "tag_name": release_tag,
            "name": release_name,
            "body": (
                "JAMA Network Open makalesi iÃ§in oluÅŸturulan gÃ¶rsel Ã¶zet.\n\n"
                f"Makale: {safe_title}\n"
                f"OluÅŸturulma tarihi: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            ),
            "draft": False,
            "prerelease": False,
        }
        response = requests.post(f"{api_base}/releases", json=release_data, headers=headers_json)
        if response.status_code != 201:
            # Ã–zel durum: Repository is empty -> README oluÅŸturarak baÅŸlatmayÄ± dene, sonra tekrar dene
            if response.status_code == 422 and "Repository is empty" in response.text:
                print("Repo boÅŸ. README.md oluÅŸturarak repo'yu baÅŸlatmayÄ± deniyoruz...")
                repo_meta = requests.get(api_base, headers=headers_json)
                default_branch = "main"
                if repo_meta.status_code == 200:
                    default_branch = repo_meta.json().get("default_branch") or "main"

                readme_content = (
                    f"# {repo_name}\n\nBu repo otomatik olarak gÃ¶rsel Ã¶zet dosyalarÄ±nÄ± (PPTX) barÄ±ndÄ±rmak iÃ§in baÅŸlatÄ±ldÄ±."
                )
                create_resp = requests.put(
                    f"{api_base}/contents/README.md",
                    headers=headers_json,
                    json={
                        "message": "Initialize repository with README",
                        "content": base64.b64encode(readme_content.encode("utf-8")).decode("utf-8"),
                        "branch": default_branch,
                    },
                )
                if create_resp.status_code in (201, 200):
                    print("README.md oluÅŸturuldu. Release tekrar oluÅŸturuluyor...")
                    time.sleep(1)
                    response = requests.post(
                        f"{api_base}/releases", json=release_data, headers=headers_json
                    )
                else:
                    msg = (
                        "Repo boÅŸ ve README oluÅŸturma baÅŸarÄ±sÄ±z: "
                        f"{create_resp.status_code} {create_resp.text}"
                    )
                    print(msg)
                    return None, msg

            if response.status_code != 201:
                msg = f"Release oluÅŸturma hatasÄ±: {response.status_code} {response.text}"
                print(msg)
                return None, msg

        release_info = response.json()
        upload_url = release_info["upload_url"].split("{")[0]

        # DosyayÄ± yÃ¼kle
        print("Dosya yÃ¼kleniyor...")
        with open(filename, "rb") as f:
            binary = f.read()
        headers_upload = {
            "Authorization": f"Bearer {github_token}",
            "Accept": "application/vnd.github+json",
            "Content-Type": "application/octet-stream",
        }
        upload_resp = requests.post(
            f"{upload_url}?name={os.path.basename(filename)}",
            data=binary,
            headers=headers_upload,
        )
        if upload_resp.status_code != 201:
            msg = f"Dosya yÃ¼kleme hatasÄ±: {upload_resp.status_code} {upload_resp.text}"
            print(msg)
            return None, msg

        asset_info = upload_resp.json()
        download_url = asset_info.get("browser_download_url")
        if not download_url:
            msg = "browser_download_url bulunamadÄ±."
            print(msg)
            return None, msg
        print(f"Dosya baÅŸarÄ±yla yÃ¼klendi: {download_url}")
        return download_url, None
    except Exception as e:
        msg = f"GitHub yÃ¼kleme hatasÄ±: {e}"
        print(msg)
        return None, msg

def create_graphical_abstract_from_url(url: str) -> str:
    # GÃœNCELLEME: Hata mesajÄ±nÄ± iÅŸlemek iÃ§in gÃ¼ncellendi
    print(f"Makale ayrÄ±ÅŸtÄ±rÄ±lÄ±yor: {url}")
    parsed_data, error = parse_jama_article(url) # ArtÄ±k iki deÄŸer alÄ±yoruz
    
    if error:
        # EÄŸer parse_jama_article bir hata dÃ¶ndÃ¼rdÃ¼yse, o hatayÄ± direkt olarak kullanÄ±cÄ±ya gÃ¶ster.
        return f"HATA: Makale verileri Ã§ekilemedi. Teknik Detay: {error}"

    print("Ä°Ã§eriÄŸe gÃ¶re tematik ikon seÃ§iliyor...")
    thematic_icon_path = select_thematic_icon(
        parsed_data.get('title', ''),
        parsed_data.get('keywords', [])
    )
    
    print("PowerPoint sunumu oluÅŸturuluyor...")
    local_filename = create_presentation(parsed_data, thematic_icon_path)
    
    # Eski davranÄ±ÅŸ: Ortam deÄŸiÅŸkenlerinden repo ve token alÄ±nÄ±p yÃ¼kleme denenir
    print("Dosya GitHub'a yÃ¼kleniyor...")
    env_repo = os.getenv("GITHUB_REPO")
    env_token = os.getenv("GITHUB_TOKEN")
    download_url = None
    if env_repo and env_token:
        download_url, upload_err = upload_to_github_release(
            local_filename,
            parsed_data.get('title', 'Bilinmeyen Makale'),
            env_repo,
            env_token,
        )
    
    if download_url:
        return f"âœ… PowerPoint sunumu baÅŸarÄ±yla oluÅŸturuldu!\n\nğŸ“¥ Ä°ndirme linki: {download_url}\n\nğŸ’¡ Bu link kalÄ±cÄ±dÄ±r ve herkese aÃ§Ä±ktÄ±r."
    else:
        extra = f"\n\nDetay: {upload_err}" if env_repo and env_token else ""
        return (
            f"âœ… PowerPoint sunumu baÅŸarÄ±yla oluÅŸturuldu: {local_filename}\n\n"
            f"âš ï¸ GitHub yÃ¼kleme servisi ÅŸu anda kullanÄ±lamÄ±yor. Dosya yerel olarak kaydedildi.{extra}"
        )


def create_graphical_abstract(url: str, github_repo: str, github_token: str) -> str:
    """
    KullanÄ±cÄ±dan alÄ±nan URL, repo (kullaniciadi/repoadi) ve token ile PPTX oluÅŸturur,
    `latest-abstract` release'ine yÃ¼kler ve herkese aÃ§Ä±k indirme linkini dÃ¶ndÃ¼rÃ¼r.
    """
    print(f"Makale ayrÄ±ÅŸtÄ±rÄ±lÄ±yor: {url}")
    parsed_data, error = parse_jama_article(url)
    if error:
        return f"HATA: Makale verileri Ã§ekilemedi. Teknik Detay: {error}"

    print("Ä°Ã§eriÄŸe gÃ¶re tematik ikon seÃ§iliyor...")
    thematic_icon_path = select_thematic_icon(
        parsed_data.get("title", ""), parsed_data.get("keywords", [])
    )

    print("PowerPoint sunumu oluÅŸturuluyor...")
    local_filename = create_presentation(parsed_data, thematic_icon_path)

    print("GitHub release oluÅŸturuluyor ve dosya yÃ¼kleniyor...")
    download_url, upload_err = upload_to_github_release(
        local_filename,
        parsed_data.get("title", "Bilinmeyen Makale"),
        github_repo,
        github_token,
    )

    if download_url:
        return (
            "âœ… PowerPoint sunumu baÅŸarÄ±yla oluÅŸturuldu!\n\n"
            f"ğŸ“¥ Ä°ndirme linki: {download_url}\n\n"
            "ğŸ’¡ Bu link kalÄ±cÄ±dÄ±r ve herkese aÃ§Ä±ktÄ±r."
        )
    return (
        f"âœ… PowerPoint sunumu baÅŸarÄ±yla oluÅŸturuldu: {local_filename}\n\n"
        f"âš ï¸ GitHub yÃ¼kleme baÅŸarÄ±sÄ±z oldu. Repo adÄ±nÄ± ve token'Ä± kontrol edin. Detay: {upload_err}"
    )

def extract_population_info(data):
    """Makale verilerinden popÃ¼lasyon bilgisini Ã§Ä±karÄ±r"""
    # Ã–nce abstract'tan dene
    if "abstract" in data and "design, setting, and participants" in data["abstract"]:
        text = data["abstract"]["design, setting, and participants"]
        
        # PopÃ¼lasyon bilgisini Ã§Ä±kar (Ã¶rnek: "115 Men, 224 Women")
        if "men" in text.lower() and "women" in text.lower():
            # SayÄ±larÄ± ve cinsiyet bilgisini Ã§Ä±kar
            import re
            numbers = re.findall(r'\d+', text)
            if len(numbers) >= 2:
                return f"{numbers[0]} Men, {numbers[1]} Women"
        
        # Alternatif format: "X participants" veya "X patients"
        elif "participants" in text.lower() or "patients" in text.lower():
            import re
            numbers = re.findall(r'\d+', text)
            if numbers:
                return f"{numbers[0]} Participants"
        
        # YaÅŸ bilgisi varsa ekle
        elif "mean age" in text.lower() or "age" in text.lower():
            import re
            age_match = re.search(r'mean age[,\s]*(\d+\.?\d*)', text.lower())
            if age_match:
                age = age_match.group(1)
                return f"Mean age: {age} years"
    
    # Fallback
    return data["abstract"].get("design, setting, and participants", "PopÃ¼lasyon bilgisi bulunamadÄ±.")

def extract_intervention_info(data):
    """Makale verilerinden mÃ¼dahale bilgisini Ã§Ä±karÄ±r"""
    if "abstract" in data and "interventions" in data["abstract"]:
        text = data["abstract"]["interventions"]
        # KatÄ±lÄ±mcÄ± sayÄ±sÄ±nÄ± Ã§Ä±kar
        import re
        numbers = re.findall(r'\d+', text)
        if numbers:
            return f"{numbers[0]} Participants analyzed"
    
    return data["abstract"].get("interventions", "MÃ¼dahale bilgisi bulunamadÄ±.")

def extract_findings_info(data):
    """Makale verilerinden bulgular bilgisini Ã§Ä±karÄ±r"""
    if "abstract" in data and "results" in data["abstract"]:
        return data["abstract"]["results"]
    
    return "Bulgular bulunamadÄ±."

def extract_settings_info(data):
    """Makale verilerinden ayarlar ve konum bilgisini Ã§Ä±karÄ±r"""
    if "abstract" in data and "design, setting, and participants" in data["abstract"]:
        text = data["abstract"]["design, setting, and participants"]
        # Konum bilgisini Ã§Ä±kar
        if "units" in text.lower() or "centers" in text.lower():
            import re
            numbers = re.findall(r'\d+', text)
            if numbers:
                return f"{numbers[0]} Psychiatric inpatient units across the US"
    
    return data["abstract"].get("design, setting, and participants", "Ayarlar ve konumlar bulunamadÄ±.")

def extract_primary_outcome_info(data):
    """Makale verilerinden birincil sonuÃ§ bilgisini Ã§Ä±karÄ±r"""
    if "abstract" in data and "main outcomes and measures" in data["abstract"]:
        return data["abstract"]["main outcomes and measures"]
    
    return "Birincil sonuÃ§ bulunamadÄ±."