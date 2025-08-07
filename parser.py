import json
import re
from pprint import pprint
from bs4 import BeautifulSoup

# Selenium kütüphaneleri
from selenium import webdriver
from selenium.webdriver.chrome.service import Service as ChromeService
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.chrome.options import Options

# PowerPoint kütüphaneleri
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# --- İKON EŞLEŞTİRME HARİTASI ---
ICON_MAP = {
    "cardiology.png": ["heart", "cardiac", "cardiology", "myocardial", "arrhythmia", "heart failure"],
    "neurology.png": ["brain", "neuro", "neurology", "stroke", "alzheimer", "parkinson", "epilepsy"],
    "oncology.png": ["cancer", "oncology", "tumor", "chemotherapy", "carcinoma"],
    "public_health.png": ["population", "public health", "mortality", "epidemiology", "opioid", "addiction"],
    "genetics.png": ["gene", "genetic", "dna", "genome", "genomics"],
}

def select_thematic_icon(article_title, article_keywords):
    search_text = (article_title.lower() + ' ' + ' '.join(article_keywords)).lower()
    for icon_file, keywords in ICON_MAP.items():
        for keyword in keywords:
            if keyword in search_text:
                print(f"'{keyword}' anahtar kelimesi bulundu. '{icon_file}' ikonu seçildi.")
                return f"icons/{icon_file}" # Klasör yoluyla birlikte döndür
    print("Tematik anahtar kelime bulunamadı. Varsayılan ikon kullanılacak.")
    return "icons/default.png" # Varsayılan ikonun da klasörde olduğunu varsayalım

def parse_jama_article(url):
    # Bu fonksiyon önceki mesajdakiyle aynı, değişiklik yok.
    # ... (önceki mesajdaki parse_jama_article fonksiyonunun tamamı buraya gelecek)
    html_content = None
    chrome_options = Options()
    chrome_options.add_argument("--headless")
    chrome_options.add_argument("--log-level=3")
    chrome_options.add_argument("user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/109.0.0.0 Safari/537.36")
    try:
        service = ChromeService(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=chrome_options)
        print("Sayfa Selenium ile açılıyor, bu işlem biraz sürebilir...")
        driver.get(url)
        html_content = driver.page_source
    except Exception as e:
        print(f"Selenium ile sayfa yüklenirken bir hata oluştu: {e}")
        return None
    finally:
        if 'driver' in locals():
            driver.quit()
    if not html_content:
        print("HTML içerik alınamadı.")
        return None
    soup = BeautifulSoup(html_content, 'html.parser')
    article_data = {
        'url': url, 'title': None, 'authors': [], 'publication_date': None,
        'doi': None, 'key_points': {}, 'abstract': {}, 'keywords': [],
        'full_text': {}, 'references': []
    }
    try:
        for script_tag in soup.find_all('script', {'type': 'application/ld+json'}):
            if script_tag.string:
                json_content = json.loads(script_tag.string)
                if json_content.get('@type') == 'MedicalScholarlyArticle':
                    article_data['publication_date'] = json_content.get('datePublished')
                    article_data['keywords'] = json_content.get('keyWords', '').split(', ')
                    break
    except (json.JSONDecodeError, AttributeError, TypeError):
        print("JSON-LD verisi ayrıştırılırken bir sorun oluştu, devam ediliyor.")
    title_tag = soup.find('h1', class_='meta-article-title')
    if title_tag:
        article_data['title'] = title_tag.get_text(strip=True)
    authors_limited = soup.select('.meta-authors--limited .wi-fullname')
    authors_remaining = soup.select('.meta-authors--remaining .wi-fullname')
    all_authors = authors_limited + authors_remaining
    author_affiliations_map = {}
    affiliations_list = soup.select('.meta-author-affiliations li')
    for li in affiliations_list:
        sup_tag = li.find('sup')
        if sup_tag:
            sup_num = sup_tag.get_text(strip=True)
            affiliation_text_div = li.find('div', class_='meta-author-name')
            if affiliation_text_div:
                affiliation_text = affiliation_text_div.get_text(strip=True).replace(sup_num, '', 1).strip()
                author_affiliations_map[sup_num] = affiliation_text
    for author in all_authors:
        name_tag = author.find('a') if author.find('a') else author
        sup_tags = name_tag.find_all('sup')
        aff_keys = [sup.get_text(strip=True) for sup in sup_tags]
        for sup in name_tag.find_all('sup'):
            sup.decompose()
        name = name_tag.get_text(strip=True).replace(',', '')
        author_info = {
            'name': name,
            'affiliations': [author_affiliations_map.get(key) for key in aff_keys if author_affiliations_map.get(key)]
        }
        article_data['authors'].append(author_info)
    doi_tag = soup.find('span', class_='meta-citation-doi')
    if doi_tag:
        article_data['doi'] = doi_tag.get_text(strip=True).replace('doi:', '')
    key_points_section = soup.find('span', class_='heading-text', string='Key Points')
    if key_points_section:
        current_element = key_points_section.find_next_sibling('p')
        while current_element and current_element.name == 'p':
            strong_tag = current_element.find('strong')
            if strong_tag:
                key = strong_tag.get_text(strip=True).lower()
                value = current_element.find('span').get_text(strip=True) if current_element.find('span') else ''
                article_data['key_points'][key] = value
            current_element = current_element.find_next_sibling()
    abstract_section = soup.find('div', id='AbstractSection')
    if abstract_section:
        for p_tag in abstract_section.find_all('p'):
            strong_tag = p_tag.find('strong')
            if strong_tag:
                key = strong_tag.get_text(strip=True).lower().replace(':', '')
                span_tag = p_tag.find('span')
                if span_tag:
                    article_data['abstract'][key] = span_tag.get_text(strip=True)
    full_text_div = soup.find('div', class_='article-full-text')
    if full_text_div:
        first_header = full_text_div.find('div', class_='section-type-section')
        if first_header:
            current_section_title = "introduction"
            article_data['full_text'][current_section_title] = ""
            for element in first_header.find_all_previous('p', class_='para'):
                 article_data['full_text'][current_section_title] = element.get_text(" ", strip=True) + "\n" + article_data['full_text'][current_section_title]
            
            for element in first_header.find_next_siblings():
                if element.name == 'div' and 'section-type-section' in element.get('class', []):
                    current_section_title = element.get_text(strip=True).lower()
                    article_data['full_text'][current_section_title] = ""
                elif element.name == 'p' and 'para' in element.get('class', []):
                    if current_section_title:
                        article_data['full_text'][current_section_title] += element.get_text(" ", strip=True) + "\n"
    references_div = soup.find('div', class_='references')
    if references_div:
        for ref in references_div.find_all('div', class_='reference'):
            ref_content = ref.find('div', class_='reference-content')
            if ref_content:
                article_data['references'].append(ref_content.get_text(" ", strip=True))
    return article_data

def create_presentation(data, icon_path):
    """
    Verilen verilerle bir PowerPoint sunumu oluşturur ve kaydeder.
    Metin kaydırma, otomatik sığdırma ve dinamik başlık boyutu eklendi.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Metin kutusu eklemek için yardımcı fonksiyon (autofit eklendi)
    def add_textbox(text, left, top, width, height, font_size=12, is_bold=False, font_color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        
        # YENİ: Metin kutuya sığmazsa fontu otomatik küçült
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True # Word wrap hala önemli
        
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.color.rgb = font_color
        p.alignment = align
        return textbox

    # Arka plan kutusu eklemek için yardımcı fonksiyon
    def add_background_box(left, top, width, height, color=RGBColor(0xF2, 0xF2, 0xF2)):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # --- Slayt Öğelerini Yerleştirme ---
    
    # 1. Logo
    try:
        slide.shapes.add_picture('jama_logo.png', Inches(0.2), Inches(0.2), height=Inches(0.4))
    except FileNotFoundError:
        print("UYARI: 'jama_logo.png' bulunamadı.")

    # 2. Ana Başlık (Dinamik Font Boyutu ile)
    title = data.get('title', 'Başlık Bulunamadı')
    # YENİ: Başlık uzunsa fontu küçült
    title_font_size = 18
    if len(title) > 100:
        title_font_size = 14
        print(f"Başlık uzun ({len(title)} karakter), font boyutu {title_font_size}pt olarak ayarlandı.")
        
    add_textbox(f"RCT: {title}", 0.3, 0.7, 9.4, 0.6, font_size=title_font_size, is_bold=True, font_color=RGBColor(0xED, 0x09, 0x73))

    # --- İçerik Kutuları ---
    # Not: Kutuların 'top' değeri (dikey konumu) 1.5'ten 1.4'e çekildi, başlığa daha çok yer bırakmak için.
    
    # 3. Nüfus (Population) Kutusu
    add_background_box(0.3, 1.5, 2.9, 3.6) # Arka plan
    add_textbox("POPULATION", 0.5, 1.6, 2.5, 0.3, font_size=11, is_bold=True, font_color=RGBColor(0xED, 0x09, 0x73))
    try:
        slide.shapes.add_picture(icon_path, Inches(0.5), Inches(2.0), height=Inches(0.8))
    except FileNotFoundError:
        print(f"UYARI: '{icon_path}' bulunamadı.")
    
    population_text = data['abstract'].get('design, setting, and participants', 'Nüfus verisi bulunamadı.')
    add_textbox(population_text, 0.5, 2.9, 2.5, 2.0, font_size=10)

    # 4. Bulgular (Findings) Kutusu
    add_background_box(3.4, 1.5, 6.3, 3.6) # Arka plan
    add_textbox("FINDINGS", 3.6, 1.6, 3.5, 0.3, font_size=11, is_bold=True, font_color=RGBColor(0xED, 0x09, 0x73))
    findings_text = data['abstract'].get('conclusions and relevance', 'Bulgular bulunamadı.')
    add_textbox(findings_text, 3.6, 2.0, 6.0, 3.0, font_size=10)
    
    # 5. Künye (Footer)
    first_author = data['authors'][0]['name'] if data['authors'] else 'Yazar Yok'
    citation = f"{first_author.split(',')[0]} L, et al. {title}; a randomized clinical trial. JAMA Netw Open. {data.get('publication_date', 'Tarih Yok')} doi:{data.get('doi', 'DOI Yok')}"
    add_textbox(citation, 0.3, 5.2, 9.4, 0.3, font_size=8)
    
    filename = "JAMA_Graphical_Abstract.pptx"
    prs.save(filename)
    print(f"\nSunum başarıyla güncellendi ve kaydedildi: {filename}")


# --- ANA ÇALIŞTIRMA BLOĞU ---
if __name__ == '__main__':
    test_url = "https://jamanetwork.com/journals/jamanetworkopen/fullarticle/2837321"
    
    print(f"Makale ayrıştırılıyor: {test_url}")
    
    # 1. Adım: Veriyi çek ve ayrıştır
    parsed_data = parse_jama_article(test_url)
    
    if parsed_data:
        print("\n--- Ayrıştırılan Makale Verileri ---")
        pprint({k: v for k, v in parsed_data.items() if k not in ['full_text', 'references']})

        # 2. Adım: İçeriğe göre ikon seç
        print("\n--- İkon Seçimi ---")
        thematic_icon_path = select_thematic_icon(
            parsed_data.get('title', ''),
            parsed_data.get('keywords', [])
        )
        
        # 3. Adım: PowerPoint sunumunu oluştur
        print("\n--- Sunum Oluşturma ---")
        create_presentation(parsed_data, thematic_icon_path)